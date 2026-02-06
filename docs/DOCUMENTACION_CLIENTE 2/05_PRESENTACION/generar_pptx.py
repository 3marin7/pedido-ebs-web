#!/usr/bin/env python3
"""
Script para generar presentaci√≥n PowerPoint profesional
Sistema de Gesti√≥n Comercial Distribuciones EBS
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def crear_presentacion():
    """Crea presentaci√≥n PowerPoint con dise√±o profesional"""
    
    # Crear presentaci√≥n
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colores
    COLOR_PRINCIPAL = RGBColor(0, 102, 204)  # Azul profesional
    COLOR_SECUNDARIO = RGBColor(255, 153, 0)  # Naranja
    COLOR_TEXTO = RGBColor(51, 51, 51)  # Gris oscuro
    COLOR_BLANCO = RGBColor(255, 255, 255)
    
    def agregar_titulo_subtitulo(slide, titulo, subtitulo=""):
        """Agrega t√≠tulo y subt√≠tulo a diapositiva"""
        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(54)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRINCIPAL
        
        if subtitulo:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitulo
            subtitle.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXTO
    
    def agregar_contenido(slide, titulo, puntos):
        """Agrega contenido con vi√±etas"""
        # T√≠tulo
        title = slide.shapes.title
        title.text = titulo
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRINCIPAL
        
        # Contenido
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, punto in enumerate(puntos):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = punto
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXTO
            p.space_before = Pt(12)
            p.space_after = Pt(12)
    
    # ========== SLIDE 1: PORTADA ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRINCIPAL
    
    # T√≠tulo principal
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Sistema de Gesti√≥n Comercial"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subt√≠tulo
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Distribuciones EBS"
    subtitle_frame.paragraphs[0].font.size = Pt(48)
    subtitle_frame.paragraphs[0].font.bold = True
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SECUNDARIO
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Pie de p√°gina
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Digitalizaci√≥n para San Andresitos | 2 de febrero de 2026"
    footer_frame.paragraphs[0].font.size = Pt(16)
    footer_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ========== SLIDE 2: CONTEXTO DEL PROBLEMA ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide2, "El Problema Actual", [
        "‚ùå Facturaci√≥n manual con errores frecuentes",
        "‚ùå P√©rdida de facturas f√≠sicas (15% mensuales)",
        "‚ùå Cobros no sistematizados sin seguimiento",
        "‚ùå Inventario desactualizado",
        "‚ùå Cero respaldos de informaci√≥n cr√≠tica"
    ])
    
    # ========== SLIDE 3: IMPACTO FINANCIERO ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide3, "Impacto Financiero Actual", [
        "üí∞ P√©rdidas mensuales: $6.3 MILLONES COP",
        "",
        "üìâ Desglose:",
        "   ‚Ä¢ Facturas extraviadas: $2.5M/mes",
        "   ‚Ä¢ Morosidad en cr√©ditos: $3M/mes",
        "   ‚Ä¢ Errores de c√°lculo: $800K/mes"
    ])
    
    # ========== SLIDE 4: OBJETIVO ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide4, "Objetivo del Proyecto", [
        "‚úÖ Digitalizar facturaci√≥n electr√≥nica",
        "‚úÖ Automatizar seguimiento de cobros",
        "‚úÖ Sincronizar inventario en tiempo real",
        "‚úÖ Generar respaldos autom√°ticos",
        "‚úÖ Permitir acceso remoto seguro 24/7"
    ])
    
    # ========== SLIDE 5: SOLUCI√ìN EBS (M√ìDULOS) ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide5, "Soluci√≥n: 6 M√≥dulos Integrados", [
        "1Ô∏è‚É£ Facturaci√≥n Electr√≥nica con numeraci√≥n √∫nica",
        "2Ô∏è‚É£ Gesti√≥n de Cuentas de Cobro",
        "3Ô∏è‚É£ Control de Inventario en Tiempo Real",
        "4Ô∏è‚É£ Cat√°logo Digital de Productos",
        "5Ô∏è‚É£ Base de Datos de Clientes",
        "6Ô∏è‚É£ Reportes Ejecutivos"
    ])
    
    # ========== SLIDE 6: BENEFICIOS TANGIBLES ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide6, "Beneficios Cuantificables", [
        "üéØ 90% menos errores en facturaci√≥n",
        "üéØ 40% reducci√≥n en morosidad",
        "üéØ 100% de facturas respaldadas",
        "üéØ 80% menos tiempo en gesti√≥n",
        "üéØ ROI: Retorno de inversi√≥n en 1.85 meses"
    ])
    
    # ========== SLIDE 7: ARQUITECTURA T√âCNICA ==========
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide7, "Arquitectura T√©cnica", [
        "üîß Frontend: React.js + Dise√±o responsivo",
        "üîß Backend: Supabase (PostgreSQL)",
        "üîß Hosting: Vercel (alta disponibilidad)",
        "üîß Seguridad: SSL/TLS + Encriptaci√≥n",
        "üîß Respaldos: Backup semanal autom√°tico"
    ])
    
    # ========== SLIDE 8: CRONOGRAMA ==========
    slide8 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide8, "Cronograma de Implementaci√≥n", [
        "üìÖ Fase 1 (Semanas 1-2): An√°lisis y Preparaci√≥n",
        "üìÖ Fase 2 (Semanas 3-4): Implementaci√≥n y Migraci√≥n",
        "üìÖ Fase 3 (Semanas 5-6): Capacitaci√≥n y Pruebas",
        "üìÖ Fase 4 (Semanas 7-8): Puesta en Marcha",
        "",
        "‚è±Ô∏è TIEMPO TOTAL: 8 semanas"
    ])
    
    # ========== SLIDE 9: INVERSI√ìN ==========
    slide9 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide9, "Inversi√≥n Requerida", [
        "üíµ Implementaci√≥n (pago √∫nico): $10.000.000 COP",
        "   ‚Üí 50% a la firma, 50% al finalizar",
        "",
        "üíµ Mensualidad (soporte completo): $400.000 COP",
        "   ‚Üí Hosting, backup, soporte, mantenimiento",
        "",
        "üìä Punto de equilibrio: 1.85 meses"
    ])
    
    # ========== SLIDE 10: COMPARATIVA ANTES/DESPU√âS ==========
    slide10 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide10, "Comparativa: Antes vs Despu√©s", [
        "ANTES:",
        "   ‚Ä¢ Facturas manuales: 15-20 min/factura",
        "   ‚Ä¢ P√©rdidas: $6.3M/mes",
        "",
        "DESPU√âS:",
        "   ‚Ä¢ Facturaci√≥n autom√°tica: 2-3 min/factura",
        "   ‚Ä¢ Recuperaci√≥n: $5.4M/mes"
    ])
    
    # ========== SLIDE 11: PR√ìXIMOS PASOS ==========
    slide11 = prs.slides.add_slide(prs.slide_layouts[5])
    agregar_contenido(slide11, "Pr√≥ximos Pasos", [
        "1Ô∏è‚É£ Revisi√≥n de documentaci√≥n (30 minutos)",
        "2Ô∏è‚É£ Reuni√≥n de aclaraci√≥n (esta semana)",
        "3Ô∏è‚É£ Demo del sistema en vivo",
        "4Ô∏è‚É£ Aprobaci√≥n y firma de contrato",
        "5Ô∏è‚É£ Inicio de implementaci√≥n"
    ])
    
    # ========== SLIDE 12: CIERRE ==========
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide12.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRINCIPAL
    
    closing_box = slide12.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    closing_frame = closing_box.text_frame
    closing_frame.text = "¬øPreguntas?"
    closing_frame.paragraphs[0].font.size = Pt(60)
    closing_frame.paragraphs[0].font.bold = True
    closing_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    closing_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    contact_box = slide12.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Edwin Mar√≠n\nDistribuciones EBS\nüìß contacto@distribucionesebs.com"
    contact_frame.paragraphs[0].font.size = Pt(20)
    contact_frame.paragraphs[0].font.color.rgb = COLOR_BLANCO
    contact_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Guardar presentaci√≥n
    output_path = "/Users/edwinmarin/pedido-ebs-web/DOCUMENTACION_CLIENTE/05_PRESENTACION/Presentacion_EBS.pptx"
    prs.save(output_path)
    print(f"‚úÖ Presentaci√≥n creada: {output_path}")
    return output_path

if __name__ == "__main__":
    crear_presentacion()
